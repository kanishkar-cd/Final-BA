from __future__ import annotations

import logging
from pathlib import Path
from app.parsers.base_parser import BaseParser

logger = logging.getLogger(__name__)


class PPTParser(BaseParser):
    """Parser for PowerPoint presentation files (.pptx, .ppt)."""

    async def parse(self, file_path: str) -> str:
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"PowerPoint document not found: {path}")

        extracted_slides: list[str] = []

        try:
            import pptx
            prs = pptx.Presentation(str(path))
            for i, slide in enumerate(prs.slides, 1):
                extracted_slides.append(f"=== Slide {i} ===")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        extracted_slides.append(shape.text.strip())
            return "\n".join(extracted_slides)
        except Exception as exc:
            logger.warning("pptx parsing failed for %s, using fallback reader: %s", path, exc)
            try:
                text_content = path.read_text(encoding="utf-8", errors="ignore")
                return text_content if text_content.strip() else f"PowerPoint Presentation: {path.name}"
            except Exception:
                return f"PowerPoint Presentation: {path.name}"
